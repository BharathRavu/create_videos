# slide_maker.py
from pptx import Presentation
import json

def make_slides(script_json, output_pptx):
    prs = Presentation()
    slides_data = json.loads(script_json)
    for slide_info in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # title+bullets
        slide.shapes.title.text = slide_info["title"]
        body = slide.shapes.placeholders[1].text_frame
        for b in slide_info["bullets"]:
            p = body.add_paragraph()
            p.text = b
            p.level = 1
    prs.save(output_pptx)
